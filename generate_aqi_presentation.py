from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

OUTPUT = Path(__file__).with_name("Air_Quality_Prediction_Project_Presentation.pptx")

BG = RGBColor(11, 18, 32)
PANEL = RGBColor(18, 30, 50)
PANEL_2 = RGBColor(27, 43, 70)
TEXT = RGBColor(244, 247, 252)
MUTED = RGBColor(181, 194, 214)
ACCENT = RGBColor(91, 192, 190)
ACCENT_2 = RGBColor(244, 166, 80)
ACCENT_3 = RGBColor(233, 90, 82)
GOOD = RGBColor(95, 201, 132)
WARN = RGBColor(250, 192, 94)
BAD = RGBColor(205, 79, 79)
LIGHT = RGBColor(235, 241, 247)


def set_slide_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=20, color=TEXT,
                bold=False, align=PP_ALIGN.LEFT, font_name="Aptos", margin=0.05):
    box = slide.shapes.add_textbox(left, top, width, height)
    box.text_frame.word_wrap = True
    box.text_frame.margin_left = Inches(margin)
    box.text_frame.margin_right = Inches(margin)
    box.text_frame.margin_top = Inches(margin)
    box.text_frame.margin_bottom = Inches(margin)
    box.text_frame.vertical_anchor = MSO_ANCHOR.TOP
    p = box.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    font = run.font
    font.name = font_name
    font.size = Pt(font_size)
    font.bold = bold
    font.color.rgb = color
    return box


def add_bullets(slide, left, top, width, height, items, font_size=18, color=TEXT,
                bullet_color=ACCENT, level_spacing=0.22):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    first = True
    for item in items:
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.text = text
        p.level = level
        p.bullet = True
        p.space_after = Pt(3)
        p.line_spacing = 1.05
        if level:
            p.left_margin = Inches(level * level_spacing)
        for run in p.runs:
            run.font.name = "Aptos"
            run.font.size = Pt(font_size - (level * 1))
            run.font.color.rgb = color
    return box


def add_header(slide, title, subtitle=None, slide_no=None):
    add_textbox(slide, Inches(0.42), Inches(0.24), Inches(8.7), Inches(0.45), title,
                font_size=24, bold=True, color=TEXT, font_name="Aptos Display")
    if subtitle:
        add_textbox(slide, Inches(0.42), Inches(0.68), Inches(8.6), Inches(0.3), subtitle,
                    font_size=11, color=MUTED)
    accent_line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.42), Inches(1.02), Inches(12.5), Inches(0.04))
    accent_line.fill.solid()
    accent_line.fill.fore_color.rgb = ACCENT
    accent_line.line.fill.background()
    if slide_no is not None:
        pill = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(12.0), Inches(0.2), Inches(0.82), Inches(0.38))
        pill.fill.solid()
        pill.fill.fore_color.rgb = PANEL_2
        pill.line.color.rgb = ACCENT
        tf = pill.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = str(slide_no)
        run.font.size = Pt(12)
        run.font.bold = True
        run.font.color.rgb = TEXT
        run.font.name = "Aptos"


def add_footer(slide, text="Air Quality Prediction System"):
    add_textbox(slide, Inches(0.42), Inches(7.02), Inches(6.5), Inches(0.22), text, font_size=9, color=MUTED)


def add_chip(slide, x, y, w, h, text, fill, text_color=TEXT, border=None):
    chip = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    chip.fill.solid()
    chip.fill.fore_color.rgb = fill
    if border is None:
        chip.line.fill.background()
    else:
        chip.line.color.rgb = border
    tf = chip.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = text_color
    return chip


def add_panel(slide, x, y, w, h, fill=PANEL, line=ACCENT):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1.25)
    return shape


def add_flow_box(slide, x, y, w, h, title, body, accent):
    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    box.fill.solid()
    box.fill.fore_color.rgb = PANEL_2
    box.line.color.rgb = accent
    box.line.width = Pt(1.5)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.06)
    tf.margin_bottom = Inches(0.03)
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = title
    r1.font.name = "Aptos Display"
    r1.font.size = Pt(13)
    r1.font.bold = True
    r1.font.color.rgb = TEXT
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = body
    r2.font.name = "Aptos"
    r2.font.size = Pt(10)
    r2.font.color.rgb = MUTED
    return box


def add_connector(slide, x1, y1, x2, y2, color=ACCENT):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    conn.line.color.rgb = color
    conn.line.width = Pt(2)
    conn.line.end_arrowhead = True
    return conn


def add_table(slide, left, top, width, height, rows, cols, data, style_fill=PANEL):
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = style_fill if r else PANEL_2
            cell.text = str(data[r][c])
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if r == 0 else PP_ALIGN.LEFT
                for run in p.runs:
                    run.font.name = "Aptos"
                    run.font.size = Pt(10 if r else 11)
                    run.font.bold = True if r == 0 else False
                    run.font.color.rgb = TEXT
    return table_shape


def add_chart(slide, left, top, width, height):
    chart_data = CategoryChartData()
    chart_data.categories = ["Delhi", "Hyderabad", "Bengaluru"]
    chart_data.add_series("Best RMSE", (73.45, 15.29, 24.08))
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        left,
        top,
        width,
        height,
        chart_data,
    ).chart
    chart.has_legend = False
    chart.value_axis.has_major_gridlines = True
    chart.value_axis.major_gridlines.format.line.color.rgb = RGBColor(70, 86, 110)
    chart.category_axis.tick_labels.font.size = Pt(10)
    chart.value_axis.tick_labels.font.size = Pt(10)
    chart.value_axis.tick_labels.font.color.rgb = TEXT
    chart.category_axis.tick_labels.font.color.rgb = TEXT
    chart.value_axis.maximum_scale = 90
    chart.value_axis.minimum_scale = 0
    plot = chart.plots[0]
    plot.has_data_labels = True
    labels = plot.data_labels
    labels.position = XL_LABEL_POSITION.OUTSIDE_END
    labels.font.size = Pt(10)
    labels.font.bold = True
    labels.font.color.rgb = TEXT
    return chart


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Slide 1: Title
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

# Decorative bands
for idx, color in enumerate([GOOD, WARN, ACCENT_2, ACCENT_3]):
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(idx * 0.16), Inches(13.333), Inches(0.16))
    band.fill.solid()
    band.fill.fore_color.rgb = color
    band.line.fill.background()

add_textbox(slide, Inches(0.7), Inches(1.1), Inches(7.5), Inches(0.7), "Air Quality Prediction System",
            font_size=30, color=TEXT, bold=True, font_name="Aptos Display")
add_textbox(slide, Inches(0.7), Inches(1.8), Inches(7.7), Inches(0.9),
            "End-to-end project presentation: data ingestion, preprocessing, ML training, Airflow orchestration, dashboarding, Docker deployment, and Raspberry Pi production validation.",
            font_size=15, color=MUTED)
add_chip(slide, Inches(0.72), Inches(2.95), Inches(1.1), Inches(0.38), "ML", PANEL_2, ACCENT)
add_chip(slide, Inches(1.9), Inches(2.95), Inches(1.2), Inches(0.38), "Airflow", PANEL_2, ACCENT_2)
add_chip(slide, Inches(3.18), Inches(2.95), Inches(1.15), Inches(0.38), "Docker", PANEL_2, GOOD)
add_chip(slide, Inches(4.42), Inches(2.95), Inches(1.65), Inches(0.38), "Raspberry Pi", PANEL_2, BAD)
add_chip(slide, Inches(6.18), Inches(2.95), Inches(1.25), Inches(0.38), "Next.js", PANEL_2, ACCENT)

right = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(8.55), Inches(1.0), Inches(4.15), Inches(4.7))
right.fill.solid()
right.fill.fore_color.rgb = PANEL
right.line.color.rgb = ACCENT
right.line.width = Pt(1.5)
rtf = right.text_frame
rtf.clear()
rtf.word_wrap = True
for i, (label, value, col) in enumerate([
    ("Project goal", "Forecast AQI and pollutants for city-level decision support", GOOD),
    ("Training stack", "Classical ML and deep learning models were benchmarked", ACCENT_2),
    ("Production stack", "Pi-hosted Docker Compose with Airflow + dashboard + PostgreSQL", ACCENT),
    ("Live interface", "Dashboard surfaces observations, predictions, and health alerts", WARN),
]):
    p = rtf.paragraphs[0] if i == 0 else rtf.add_paragraph()
    p.space_after = Pt(10)
    r = p.add_run()
    r.text = f"{label}: "
    r.font.name = "Aptos"
    r.font.size = Pt(13)
    r.font.bold = True
    r.font.color.rgb = col
    r2 = p.add_run()
    r2.text = value
    r2.font.name = "Aptos"
    r2.font.size = Pt(13)
    r2.font.color.rgb = TEXT
add_textbox(slide, Inches(8.75), Inches(5.85), Inches(3.7), Inches(0.25), "Final year project | verified on the Pi production stack", font_size=10, color=MUTED)

# Slide 2: agenda
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_header(slide, "Presentation Roadmap", "What the project does, how it was built, and how it runs in production.", 2)
add_panel(slide, Inches(0.55), Inches(1.4), Inches(12.15), Inches(5.0))
agenda = [
    "Problem statement and project objectives",
    "Data sources and preprocessing pipeline",
    "Model training, baselines, and benchmark comparisons",
    "Airflow orchestration for ingestion and forecasts",
    "Dockerized Raspberry Pi deployment",
    "Next.js dashboard, API routes, and live monitoring",
    "Verification results, lessons learned, and next steps",
]
add_bullets(slide, Inches(0.95), Inches(1.75), Inches(11.2), Inches(4.2), agenda, font_size=20)
add_footer(slide)

# Slide 3: problem
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_header(slide, "Problem Statement and Goal", "Why this project exists and what the end-user should get from it.", 3)
add_panel(slide, Inches(0.55), Inches(1.35), Inches(5.9), Inches(5.3))
add_panel(slide, Inches(6.75), Inches(1.35), Inches(5.95), Inches(5.3), fill=RGBColor(16, 27, 45), line=ACCENT_2)
add_textbox(slide, Inches(0.82), Inches(1.65), Inches(5.2), Inches(0.35), "The problem", font_size=18, color=ACCENT, bold=True)
add_bullets(slide, Inches(0.82), Inches(2.0), Inches(5.0), Inches(3.7), [
    "Air quality changes quickly and is difficult to interpret from raw sensor values alone.",
    "The data is noisy, has gaps, and comes from multiple pipelines and frequencies.",
    "A production system must combine forecasting, health context, and operator visibility.",
    "The project therefore needs both an ML model and a user-facing dashboard.",
], font_size=18)
add_textbox(slide, Inches(7.0), Inches(1.65), Inches(5.2), Inches(0.35), "Project goals", font_size=18, color=ACCENT_2, bold=True)
add_bullets(slide, Inches(7.0), Inches(2.0), Inches(5.3), Inches(3.7), [
    "Build cleaned time-series datasets for training and evaluation.",
    "Train and compare multiple forecasting models.",
    "Automate ingestion and prediction with Airflow on the Pi.",
    "Expose live observations and forecasts through a dashboard and APIs.",
    "Deploy the full stack in Docker so it can run consistently on the Raspberry Pi.",
], font_size=18)
add_footer(slide)

# Slide 4: data pipeline
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_header(slide, "Data Sources and Preprocessing", "The raw data is cleaned before any model training or dashboard use.", 4)
add_panel(slide, Inches(0.55), Inches(1.35), Inches(12.15), Inches(5.3))
add_flow_box(slide, Inches(0.82), Inches(1.82), Inches(2.0), Inches(1.15), "Raw inputs", "feeds.csv, ThingSpeak feeds, city observation tables", GOOD)
add_flow_box(slide, Inches(3.05), Inches(1.82), Inches(2.0), Inches(1.15), "Timestamp handling", "Detect datetime column, sort index, normalize to UTC", ACCENT)
add_flow_box(slide, Inches(5.28), Inches(1.82), Inches(2.0), Inches(1.15), "Cleaning", "Rename fields, coerce numeric types, drop null columns, mark outliers", WARN)
add_flow_box(slide, Inches(7.51), Inches(1.82), Inches(2.0), Inches(1.15), "Imputation", "Rolling mean, forward-fill, backward-fill for gaps", ACCENT_2)
add_flow_box(slide, Inches(9.74), Inches(1.82), Inches(2.4), Inches(1.15), "Resample", "Standardize to 15-minute or hourly series for training", BAD)
for x1, x2 in [(2.82, 3.05), (5.05, 5.28), (7.28, 7.51), (9.51, 9.74)]:
    add_connector(slide, Inches(x1), Inches(2.39), Inches(x2), Inches(2.39))
add_textbox(slide, Inches(0.95), Inches(3.4), Inches(4.0), Inches(0.3), "Preprocessing highlights", font_size=16, color=ACCENT, bold=True)
add_bullets(slide, Inches(0.95), Inches(3.72), Inches(5.2), Inches(2.2), [
    "Automatic datetime detection prevents manual column errors.",
    "Outliers are replaced with missing values before filling.",
    "Aggregation produces stable windows for downstream training.",
], font_size=17)
add_textbox(slide, Inches(6.45), Inches(3.4), Inches(4.0), Inches(0.3), "Artifacts used in the repo", font_size=16, color=ACCENT_2, bold=True)
add_bullets(slide, Inches(6.45), Inches(3.72), Inches(5.0), Inches(2.2), [
    "phase1_2_air_quality.ipynb for loading, cleaning, and resampling.",
    "feeds_cleaned_15T.csv and feeds_cleaned.csv for processed training data.",
    "Analysis notebooks for city-level QA and benchmark generation.",
], font_size=17)
add_footer(slide)

# Slide 5: modeling
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_header(slide, "Model Training and Benchmarking", "Multiple model families were compared before deciding what to deploy on the Pi.", 5)
add_panel(slide, Inches(0.55), Inches(1.35), Inches(6.2), Inches(5.3))
add_panel(slide, Inches(6.95), Inches(1.35), Inches(5.75), Inches(5.3), fill=RGBColor(16, 27, 45), line=ACCENT)
add_textbox(slide, Inches(0.82), Inches(1.65), Inches(5.5), Inches(0.3), "Training setup", font_size=18, color=ACCENT, bold=True)
add_bullets(slide, Inches(0.82), Inches(1.98), Inches(5.4), Inches(3.9), [
    "Input window: 168 hours of history; forecast horizon: 24 hours.",
    "Train/validation/test split: 70/15/15 with StandardScaler on the feature set.",
    "Baselines: persistence, Random Forest, SVR, XGBoost, LightGBM, CatBoost.",
    "Deep learning experiments: LSTM, CNN, BiLSTM, TCN, Transformer/TFT variants.",
    "Metrics: RMSE, MAE, R2, and drift-oriented evaluation such as MAPE.",
], font_size=17)
add_textbox(slide, Inches(7.22), Inches(1.65), Inches(4.9), Inches(0.3), "Why multiple families", font_size=18, color=ACCENT_2, bold=True)
add_bullets(slide, Inches(7.22), Inches(1.98), Inches(4.9), Inches(3.9), [
    "Classical ML gives strong performance with easy Pi deployment.",
    "Deep learning provides a research baseline and can capture sequence effects.",
    "The final deployment chooses models that are accurate and ARM-friendly.",
    "Saved artifacts include .joblib and .pth files plus model metadata.",
], font_size=17)
add_footer(slide)

# Slide 6: benchmark results
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_header(slide, "Benchmark Results and Model Choice", "The best model differs by city, so the deployment uses the city-specific winner.", 6)
add_panel(slide, Inches(0.55), Inches(1.35), Inches(6.0), Inches(5.3))
add_panel(slide, Inches(6.75), Inches(1.35), Inches(6.0), Inches(5.3), fill=RGBColor(16, 27, 45), line=ACCENT_2)
add_textbox(slide, Inches(0.82), Inches(1.65), Inches(5.4), Inches(0.3), "Best deployed models", font_size=18, color=ACCENT, bold=True)
bench_table = [
    ["City", "Model", "RMSE", "MAE", "R2"],
    ["Delhi", "LightGBM", "73.45", "44.85", "0.65"],
    ["Hyderabad", "XGBoost", "15.29", "10.63", "0.59"],
    ["Bengaluru", "CatBoost", "24.08", "15.52", "0.28"],
]
add_table(slide, Inches(0.8), Inches(2.05), Inches(5.45), Inches(2.2), 4, 5, bench_table)
add_textbox(slide, Inches(0.82), Inches(4.55), Inches(5.2), Inches(0.26), "Production rationale", font_size=16, color=ACCENT_2, bold=True)
add_bullets(slide, Inches(0.82), Inches(4.84), Inches(5.25), Inches(1.35), [
    "Classical ML models are fast to load and run on Raspberry Pi.",
    "They avoid CUDA and heavyweight PyTorch dependencies in production.",
], font_size=16)
add_textbox(slide, Inches(7.0), Inches(1.65), Inches(5.2), Inches(0.3), "RMSE comparison", font_size=18, color=ACCENT_2, bold=True)
add_chart(slide, Inches(7.05), Inches(2.05), Inches(5.2), Inches(3.7))
add_footer(slide)

# Slide 7: deployment
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_header(slide, "Production Deployment on Raspberry Pi", "The live stack runs in Docker Compose on the Pi and serves both Airflow and the dashboard.", 7)
add_panel(slide, Inches(0.55), Inches(1.35), Inches(12.15), Inches(5.3))
add_textbox(slide, Inches(0.82), Inches(1.65), Inches(5.0), Inches(0.3), "Services in the Pi stack", font_size=18, color=ACCENT, bold=True)
svc_table = [
    ["Service", "Purpose", "Port"],
    ["PostgreSQL", "Stores observations and forecast data", "5432"],
    ["Airflow webserver", "Shows DAG status and run history", "8080"],
    ["Airflow scheduler / triggerer", "Runs ingestion and forecasting jobs", "internal"],
    ["Dashboard", "Serves the Next.js UI and APIs", "3000"],
]
add_table(slide, Inches(0.8), Inches(2.05), Inches(6.6), Inches(2.6), 5, 3, svc_table)
add_textbox(slide, Inches(7.8), Inches(1.65), Inches(4.4), Inches(0.3), "Why Docker matters", font_size=18, color=ACCENT_2, bold=True)
add_bullets(slide, Inches(7.8), Inches(2.0), Inches(4.2), Inches(2.7), [
    "The same compose file brings up the whole stack consistently.",
    "Repo-root mounts keep the code and data paths aligned on the Pi.",
    "Health checks and service ordering reduce startup issues.",
    "The dashboard builds and starts automatically in the container.",
], font_size=17)
add_chip(slide, Inches(7.82), Inches(5.05), Inches(1.45), Inches(0.38), "Pi verified", PANEL_2, GOOD)
add_chip(slide, Inches(9.38), Inches(5.05), Inches(1.6), Inches(0.38), "Airflow healthy", PANEL_2, ACCENT)
add_chip(slide, Inches(11.1), Inches(5.05), Inches(1.2), Inches(0.38), "Dashboard live", PANEL_2, ACCENT_2)
add_footer(slide)

# Slide 8: airflow
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_header(slide, "Airflow Orchestration", "Airflow handles ingestion, forecasting, and refresh cycles so the system stays automated.", 8)
add_panel(slide, Inches(0.55), Inches(1.35), Inches(12.15), Inches(5.3))
add_flow_box(slide, Inches(0.82), Inches(2.0), Inches(2.05), Inches(1.15), "Ingest", "Open-Meteo / city data fetched on schedule", GOOD)
add_flow_box(slide, Inches(3.12), Inches(2.0), Inches(2.05), Inches(1.15), "Normalize", "Timestamp handling, upserts, archival, and watermark checks", ACCENT)
add_flow_box(slide, Inches(5.42), Inches(2.0), Inches(2.05), Inches(1.15), "Forecast", "Load latest window and generate city-specific predictions", WARN)
add_flow_box(slide, Inches(7.72), Inches(2.0), Inches(2.05), Inches(1.15), "Store", "Persist outputs for dashboard/API consumption", ACCENT_2)
add_flow_box(slide, Inches(10.02), Inches(2.0), Inches(2.05), Inches(1.15), "Monitor", "Run health checks and trigger retraining when drift grows", BAD)
for x1, x2 in [(2.87, 3.12), (5.17, 5.42), (7.47, 7.72), (9.77, 10.02)]:
    add_connector(slide, Inches(x1), Inches(2.57), Inches(x2), Inches(2.57))
add_textbox(slide, Inches(0.88), Inches(3.65), Inches(5.3), Inches(0.3), "Key DAG behavior", font_size=18, color=ACCENT_2, bold=True)
add_bullets(slide, Inches(0.88), Inches(3.98), Inches(5.3), Inches(1.85), [
    "The forecast DAG runs hourly at 15 minutes past the hour.",
    "Tasks are separated into generation and loading steps.",
    "Retries and logging are built in for operational stability.",
], font_size=16)
add_textbox(slide, Inches(6.75), Inches(3.65), Inches(5.4), Inches(0.3), "Files on disk", font_size=18, color=ACCENT, bold=True)
add_bullets(slide, Inches(6.75), Inches(3.98), Inches(5.4), Inches(1.85), [
    "aq_forecast_dag.py defines the hourly forecasting workflow.",
    "deployment/pi_airflow/docker-compose.yml brings the stack up on the Pi.",
    "deployment_models/* contain the trained city models and metadata.",
], font_size=16)
add_footer(slide)

# Slide 9: dashboard and APIs
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_header(slide, "Dashboard and API Layer", "The Next.js app turns the ML output into usable information for people and operators.", 9)
add_panel(slide, Inches(0.55), Inches(1.35), Inches(6.05), Inches(5.3))
add_panel(slide, Inches(6.8), Inches(1.35), Inches(5.9), Inches(5.3), fill=RGBColor(16, 27, 45), line=ACCENT)
add_textbox(slide, Inches(0.82), Inches(1.65), Inches(5.0), Inches(0.3), "User-facing pages", font_size=18, color=ACCENT, bold=True)
add_bullets(slide, Inches(0.82), Inches(2.0), Inches(5.2), Inches(2.5), [
    "/dashboard/cities for the ranking grid.",
    "/dashboard/city/[citySlug] for a single-city view.",
    "/air-quality for AQI analysis and city comparisons.",
], font_size=17)
add_textbox(slide, Inches(0.82), Inches(4.0), Inches(5.0), Inches(0.3), "What the UI shows", font_size=18, color=ACCENT_2, bold=True)
add_bullets(slide, Inches(0.82), Inches(4.32), Inches(5.2), Inches(1.65), [
    "AQI category, health advisory, hourly trend, and pollutant breakdown.",
    "Forecast timeline and model summary for the selected city.",
], font_size=17)
add_textbox(slide, Inches(7.08), Inches(1.65), Inches(5.0), Inches(0.3), "API routes", font_size=18, color=ACCENT_2, bold=True)
add_bullets(slide, Inches(7.08), Inches(2.0), Inches(5.0), Inches(2.8), [
    "/api/status",
    "/api/observations",
    "/api/observations/[citySlug]",
    "/api/predictions/[citySlug]",
], font_size=18)
add_textbox(slide, Inches(7.08), Inches(4.05), Inches(5.0), Inches(0.3), "Why the route fixes mattered", font_size=18, color=ACCENT, bold=True)
add_bullets(slide, Inches(7.08), Inches(4.37), Inches(5.0), Inches(1.55), [
    "City slugs are normalized to warehouse slugs before querying the DB.",
    "That keeps the API aligned with the actual data model and the Pi deployment.",
], font_size=16)
add_footer(slide)

# Slide 10: architecture diagram
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_header(slide, "End-to-End Architecture", "From raw data to forecasts to the production dashboard.", 10)
add_panel(slide, Inches(0.55), Inches(1.35), Inches(12.15), Inches(5.3))
add_flow_box(slide, Inches(0.82), Inches(2.0), Inches(1.95), Inches(1.1), "Sensors", "ThingSpeak / Open-Meteo / feeds.csv", GOOD)
add_flow_box(slide, Inches(2.98), Inches(2.0), Inches(1.95), Inches(1.1), "Airflow", "Schedules, retries, and ETL jobs", ACCENT)
add_flow_box(slide, Inches(5.14), Inches(2.0), Inches(1.95), Inches(1.1), "PostgreSQL", "Observations + forecast tables", WARN)
add_flow_box(slide, Inches(7.3), Inches(2.0), Inches(1.95), Inches(1.1), "ML models", "LightGBM, XGBoost, CatBoost, DL baselines", ACCENT_2)
add_flow_box(slide, Inches(9.46), Inches(2.0), Inches(1.95), Inches(1.1), "Next.js", "APIs + dashboard pages", BAD)
add_flow_box(slide, Inches(11.62), Inches(2.0), Inches(0.78), Inches(1.1), "User", "Live view", GOOD)
for x1, x2 in [(2.77, 2.98), (4.93, 5.14), (7.09, 7.3), (9.25, 9.46), (11.41, 11.62)]:
    add_connector(slide, Inches(x1), Inches(2.55), Inches(x2), Inches(2.55))
add_textbox(slide, Inches(0.95), Inches(3.85), Inches(11.2), Inches(1.6),
            "The architecture separates scheduling, storage, training, and display. That makes the system easier to reason about, easier to deploy on the Pi, and safer to update without breaking the user-facing dashboard.",
            font_size=20, color=TEXT)
add_footer(slide)

# Slide 11: verification
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_header(slide, "Verification and Current Status", "What was checked after the code and deployment updates.", 11)
add_panel(slide, Inches(0.55), Inches(1.35), Inches(12.15), Inches(5.3))
add_textbox(slide, Inches(0.82), Inches(1.65), Inches(5.0), Inches(0.3), "Verified results", font_size=18, color=ACCENT, bold=True)
add_bullets(slide, Inches(0.82), Inches(2.0), Inches(5.6), Inches(3.0), [
    "Local production build completed successfully.",
    "Pi services were up: Airflow webserver, scheduler, triggerer, PostgreSQL, and dashboard.",
    "The dashboard container started and served on port 3000.",
    "Airflow health checks showed healthy scheduler and triggerer status.",
    "Live observations endpoint returned current AQI ranking data.",
], font_size=17)
add_textbox(slide, Inches(6.95), Inches(1.65), Inches(5.0), Inches(0.3), "Interpretation", font_size=18, color=ACCENT_2, bold=True)
add_bullets(slide, Inches(6.95), Inches(2.0), Inches(5.0), Inches(2.7), [
    "The stack is not just documented; it is running on the Pi.",
    "The route fixes are consistent with the deployed data model.",
    "The remaining warning in the status snapshot is about missing paths, not a broken core service.",
], font_size=17)
add_chip(slide, Inches(7.0), Inches(5.0), Inches(1.25), Inches(0.38), "Build OK", PANEL_2, GOOD)
add_chip(slide, Inches(8.35), Inches(5.0), Inches(1.6), Inches(0.38), "Airflow healthy", PANEL_2, ACCENT)
add_chip(slide, Inches(10.08), Inches(5.0), Inches(1.5), Inches(0.38), "Dashboard live", PANEL_2, ACCENT_2)
add_chip(slide, Inches(11.7), Inches(5.0), Inches(0.9), Inches(0.38), "AQ data", PANEL_2, BAD)
add_footer(slide)

# Slide 12: next steps
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_header(slide, "Lessons Learned and Next Steps", "The project is working; the remaining work is mostly product hardening and analytics depth.", 12)
add_panel(slide, Inches(0.55), Inches(1.35), Inches(5.95), Inches(5.3))
add_panel(slide, Inches(6.8), Inches(1.35), Inches(5.9), Inches(5.3), fill=RGBColor(16, 27, 45), line=ACCENT_2)
add_textbox(slide, Inches(0.82), Inches(1.65), Inches(5.0), Inches(0.3), "What worked", font_size=18, color=ACCENT, bold=True)
add_bullets(slide, Inches(0.82), Inches(2.0), Inches(5.2), Inches(2.6), [
    "City-specific model selection improved production realism.",
    "Docker and Airflow made the Pi deployment reproducible.",
    "The dashboard gives a practical view of AQI and forecasts.",
], font_size=17)
add_textbox(slide, Inches(0.82), Inches(4.2), Inches(5.0), Inches(0.3), "Next improvements", font_size=18, color=ACCENT_2, bold=True)
add_bullets(slide, Inches(0.82), Inches(4.55), Inches(5.2), Inches(1.7), [
    "Forecast history table and better accuracy tracking over time.",
    "Materialized views and caching for faster dashboard queries.",
    "More comparison, charting, and mobile-first UX improvements.",
], font_size=16)
add_textbox(slide, Inches(7.08), Inches(1.65), Inches(5.0), Inches(0.3), "Final takeaway", font_size=18, color=ACCENT_2, bold=True)
add_textbox(slide, Inches(7.08), Inches(2.1), Inches(5.0), Inches(1.7),
            "This project is a complete air-quality prediction platform: cleaned data, benchmarked models, automated forecasts, a live dashboard, and a Raspberry Pi production deployment.",
            font_size=22, color=TEXT, bold=True)
add_chip(slide, Inches(7.12), Inches(4.55), Inches(1.65), Inches(0.38), "Prediction", PANEL_2, ACCENT)
add_chip(slide, Inches(8.92), Inches(4.55), Inches(1.6), Inches(0.38), "Deployment", PANEL_2, GOOD)
add_chip(slide, Inches(10.68), Inches(4.55), Inches(1.5), Inches(0.38), "Dashboard", PANEL_2, WARN)
add_chip(slide, Inches(12.02), Inches(4.55), Inches(0.55), Inches(0.38), "Pi", PANEL_2, BAD)
add_textbox(slide, Inches(7.08), Inches(5.4), Inches(4.7), Inches(0.35), "Questions", font_size=24, bold=True, color=GOOD, font_name="Aptos Display")
add_footer(slide)

prs.save(str(OUTPUT))
print(f"Saved {OUTPUT}")
