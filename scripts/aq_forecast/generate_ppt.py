"""
generate_ppt.py — Generate a detailed, professional PowerPoint presentation
for the Air Quality Prediction Pipeline project.

Team Members:
    - Sakshith Srihari  (1602-22-748-011)
    - Kamal Bura         (1602-22-748-302)
    - Abdul Muttalib     (1602-22-748-046)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ─── Color Palette ──────────────────────────────────────────────
DARK_BG       = RGBColor(0x1A, 0x1A, 0x2E)   # Deep navy
ACCENT_BLUE   = RGBColor(0x00, 0x96, 0xD6)   # Bright blue
ACCENT_TEAL   = RGBColor(0x00, 0xBF, 0xA5)   # Teal/green
ACCENT_ORANGE = RGBColor(0xFF, 0x8A, 0x00)   # Orange
ACCENT_PURPLE = RGBColor(0x7C, 0x4D, 0xFF)   # Purple
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY    = RGBColor(0xE0, 0xE0, 0xE0)
MEDIUM_GRAY   = RGBColor(0xA0, 0xA0, 0xA0)
DARK_TEXT      = RGBColor(0x2D, 0x2D, 0x2D)
CARD_BG        = RGBColor(0x22, 0x22, 0x3A)   # Slightly lighter navy
GREEN          = RGBColor(0x4C, 0xAF, 0x50)
RED            = RGBColor(0xF4, 0x43, 0x36)
YELLOW         = RGBColor(0xFF, 0xC1, 0x07)

# ─── Helpers ────────────────────────────────────────────────────

def set_slide_bg(slide, color):
    """Set the background color of a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_rect(slide, left, top, width, height, fill_color, border_color=None):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=14,
                color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Segoe UI"):
    """Add a text box with styled text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=13,
                    color=LIGHT_GRAY, bullet_color=ACCENT_TEAL, font_name="Segoe UI"):
    """Add a bulleted text list."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        # Add bullet character
        run_bullet = p.add_run()
        run_bullet.text = "● "
        run_bullet.font.size = Pt(font_size - 2)
        run_bullet.font.color.rgb = bullet_color
        run_bullet.font.name = font_name

        run_text = p.add_run()
        run_text.text = item
        run_text.font.size = Pt(font_size)
        run_text.font.color.rgb = color
        run_text.font.name = font_name

        p.space_after = Pt(6)
    return txBox


def add_section_header(slide, text, subtitle=""):
    """Add a styled section header bar at the top."""
    # Accent line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Pt(4)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_BLUE
    line.line.fill.background()

    # Title
    add_textbox(slide, Inches(0.6), Inches(0.35), Inches(10), Inches(0.7),
                text, font_size=30, color=WHITE, bold=True)

    if subtitle:
        add_textbox(slide, Inches(0.6), Inches(1.0), Inches(10), Inches(0.4),
                    subtitle, font_size=14, color=MEDIUM_GRAY, bold=False)


def add_card(slide, left, top, width, height, title, body_items,
             accent=ACCENT_BLUE, title_size=15, body_size=12):
    """Add a card-style box with title and bullet items."""
    # Card background
    card = add_shape_rect(slide, left, top, width, height, CARD_BG, accent)

    # Accent bar at top of card
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, Pt(3)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()

    # Title
    add_textbox(slide, left + Inches(0.15), top + Inches(0.1),
                width - Inches(0.3), Inches(0.35),
                title, font_size=title_size, color=accent, bold=True)

    # Body
    if body_items:
        add_bullet_list(slide, left + Inches(0.15), top + Inches(0.45),
                        width - Inches(0.3), height - Inches(0.5),
                        body_items, font_size=body_size, color=LIGHT_GRAY,
                        bullet_color=accent)

    return card


def add_table(slide, left, top, rows, cols, col_widths, data, header_color=ACCENT_BLUE):
    """Add a styled table."""
    table_shape = slide.shapes.add_table(rows, cols, left, top,
                                         sum(col_widths), Inches(0.4) * rows)
    table = table_shape.table

    for j, w in enumerate(col_widths):
        table.columns[j].width = w

    for i in range(rows):
        for j in range(cols):
            cell = table.cell(i, j)
            cell.text = str(data[i][j])
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(11)
                paragraph.font.name = "Segoe UI"
                paragraph.alignment = PP_ALIGN.CENTER

                if i == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
                else:
                    paragraph.font.color.rgb = LIGHT_GRAY

            # Cell fill
            if i == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = CARD_BG if i % 2 == 0 else DARK_BG

    return table_shape


# ═══════════════════════════════════════════════════════════════
#  SLIDE CREATION
# ═══════════════════════════════════════════════════════════════

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # ══════════════════════════════════════════
    # SLIDE 1: TITLE SLIDE
    # ══════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_slide_bg(slide, DARK_BG)

    # Decorative top gradient bar
    for i, c in enumerate([ACCENT_BLUE, ACCENT_TEAL, ACCENT_PURPLE, ACCENT_ORANGE]):
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(i * 3.33), Inches(0), Inches(3.34), Pt(5)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = c
        bar.line.fill.background()

    # Main title
    add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(1.2),
                "Air Quality Index Prediction System",
                font_size=42, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # Subtitle
    add_textbox(slide, Inches(0.8), Inches(2.7), Inches(11.5), Inches(0.6),
                "Deep Learning & Gradient Boosting for 24-Hour AQI Forecasting — Hyderabad, Telangana",
                font_size=18, color=ACCENT_TEAL, bold=False, alignment=PP_ALIGN.CENTER)

    # Divider line
    div = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(3.6), Inches(4.33), Pt(2)
    )
    div.fill.solid()
    div.fill.fore_color.rgb = ACCENT_BLUE
    div.line.fill.background()

    # Team members
    team = [
        ("Sakshith Srihari", "1602-22-748-011"),
        ("Kamal Bura", "1602-22-748-302"),
        ("Abdul Muttalib", "1602-22-748-046"),
    ]
    for i, (name, roll) in enumerate(team):
        x = Inches(2.0 + i * 3.5)
        add_textbox(slide, x, Inches(4.0), Inches(3.2), Inches(0.4),
                    name, font_size=16, color=WHITE, bold=True,
                    alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x, Inches(4.4), Inches(3.2), Inches(0.3),
                    roll, font_size=13, color=MEDIUM_GRAY, bold=False,
                    alignment=PP_ALIGN.CENTER)

    # Footer
    add_textbox(slide, Inches(0.8), Inches(5.5), Inches(11.5), Inches(0.8),
                "B.Tech Final Year Project  •  Department of Computer Science & Engineering\n"
                "9 Model Architectures  •  PyTorch + XGBoost + LightGBM  •  Optuna HPO",
                font_size=12, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)

    # ══════════════════════════════════════════
    # SLIDE 2: TABLE OF CONTENTS
    # ══════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_section_header(slide, "Presentation Outline")

    sections = [
        ("01", "Problem Statement & Motivation", ACCENT_BLUE),
        ("02", "Literature Review", ACCENT_TEAL),
        ("03", "Data Acquisition & Sources", ACCENT_ORANGE),
        ("04", "System Architecture & Pipeline", ACCENT_PURPLE),
        ("05", "Feature Engineering", ACCENT_BLUE),
        ("06", "Deep Learning Models (7 Architectures)", ACCENT_TEAL),
        ("07", "Gradient Boosting Models (XGBoost + LightGBM)", ACCENT_ORANGE),
        ("08", "Hyperparameter Optimization (Optuna)", ACCENT_PURPLE),
        ("09", "Training Strategy & Implementation", ACCENT_BLUE),
        ("10", "Evaluation Metrics & Methodology", ACCENT_TEAL),
        ("11", "Data Analysis & Statistics", ACCENT_ORANGE),
        ("12", "Technology Stack", ACCENT_PURPLE),
        ("13", "Multi-Agent Development System", ACCENT_BLUE),
        ("14", "Results & Future Work", ACCENT_TEAL),
    ]

    for i, (num, title, color) in enumerate(sections):
        row = i // 2
        col = i % 2
        x = Inches(0.6 + col * 6.2)
        y = Inches(1.6 + row * 0.68)

        # Number badge
        badge = add_shape_rect(slide, x, y, Inches(0.5), Inches(0.45), color)
        badge.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = badge.text_frame.paragraphs[0].add_run()
        run.text = num
        run.font.size = Pt(13)
        run.font.color.rgb = WHITE
        run.font.bold = True
        run.font.name = "Segoe UI"

        add_textbox(slide, x + Inches(0.6), y + Inches(0.05), Inches(5.2), Inches(0.4),
                    title, font_size=14, color=LIGHT_GRAY)

    # ══════════════════════════════════════════
    # SLIDE 3: PROBLEM STATEMENT & MOTIVATION
    # ══════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_section_header(slide, "Problem Statement & Motivation",
                       "Why Air Quality Prediction Matters")

    # Left card: Problem
    add_card(slide, Inches(0.5), Inches(1.6), Inches(5.8), Inches(2.8),
             "THE PROBLEM", [
                 "Hyderabad ranks among India's most polluted cities",
                 "7.0 million premature deaths/year linked to air pollution (WHO)",
                 "Current AQI monitoring is reactive — no early warning system",
                 "Manual CPCB readings provide snapshots, not forecasts",
                 "Citizens lack actionable 24-hour pollution forecasts",
             ], accent=RED)

    # Right card: Our Solution
    add_card(slide, Inches(6.8), Inches(1.6), Inches(5.8), Inches(2.8),
             "OUR SOLUTION", [
                 "24-hour ahead AQI prediction using 9 ML/DL architectures",
                 "Hourly granularity from 7 CPCB monitoring stations",
                 "Automated pipeline: data → train → evaluate → compare",
                 "India NAQI standard (6 categories, 0–500 scale)",
                 "Enable proactive health advisories & policy decisions",
             ], accent=GREEN)

    # Bottom card: Key Objectives
    add_card(slide, Inches(0.5), Inches(4.8), Inches(12.1), Inches(2.0),
             "KEY OBJECTIVES", [
                 "Compare 7 deep learning architectures (RNN, LSTM, BiLSTM, Attention variants, TCN, Transformer) against 2 gradient boosting models (XGBoost, LightGBM)",
                 "Achieve R² > 0.85 for 24-hour AQI forecasting across all seasons and pollution categories",
                 "Automate hyperparameter optimization with Optuna (50 trials/model, TPE sampler, median pruning)",
                 "Build a reproducible, modular pipeline with multi-agent code review for production quality",
             ], accent=ACCENT_BLUE, body_size=12)

    # ══════════════════════════════════════════
    # SLIDE 4: LITERATURE REVIEW
    # ══════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_section_header(slide, "Literature Review",
                       "Key Related Work in AQI Prediction")

    papers = [
        ["Author(s)", "Year", "Method", "Region", "Key Finding"],
        ["Barthwal et al.", "2024", "LSTM + Attention", "India (Multiple)", "Attention mechanism improved 24h forecast accuracy by 12% over vanilla LSTM"],
        ["Sreenivasulu et al.", "2024", "BiLSTM + GRU", "Hyderabad", "Bidirectional models captured diurnal AQI cycles better; R²=0.89"],
        ["Liu et al.", "2025", "Transformer", "Beijing", "Self-attention over long sequences outperformed CNN-LSTM hybrids"],
        ["Ravindiran et al.", "2025", "XGBoost + DL", "Chennai", "Gradient boosting competitive with DL for tabular AQ features"],
        ["Ansari et al.", "2025", "TCN + Optuna", "Delhi NCR", "TCN with dilated convolutions achieved best RMSE with 3× faster training"],
        ["Kumar & Singh", "2024", "Ensemble", "India", "Model ensembles reduced variance; no single architecture dominated all metrics"],
    ]

    add_table(slide, Inches(0.5), Inches(1.6), len(papers), 5,
              [Inches(1.8), Inches(0.7), Inches(2.0), Inches(1.8), Inches(6.0)],
              papers, header_color=ACCENT_TEAL)

    # Research gap
    add_card(slide, Inches(0.5), Inches(5.3), Inches(12.1), Inches(1.5),
             "RESEARCH GAP ADDRESSED", [
                 "No existing study benchmarks all 9 architectures (RNN → Transformer + GB) on the same Hyderabad dataset with unified preprocessing",
                 "Most studies use limited date ranges; our pipeline covers 2 full years (2023–2024) of hourly data from 7 stations",
                 "Optuna-driven HPO across all models ensures fair, optimized comparison — not default hyperparameters",
             ], accent=ACCENT_PURPLE, body_size=12)

    # ══════════════════════════════════════════
    # SLIDE 5: DATA ACQUISITION
    # ══════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_section_header(slide, "Data Acquisition & Sources",
                       "Multi-Source Hourly Data for Hyderabad")

    # Left: Data sources
    add_card(slide, Inches(0.5), Inches(1.6), Inches(5.8), Inches(2.5),
             "DATA SOURCES", [
                 "Open-Meteo Air Quality API (CAMS global reanalysis)",
                 "Open-Meteo Archive API (historical weather data)",
                 "7 CPCB CAAQMS station coordinates for Hyderabad",
                 "No API key required — fully open & reproducible",
                 "2-year period: Jan 2023 – Dec 2024 (17,544 hours/station)",
             ], accent=ACCENT_BLUE)

    # Right: Stations
    stations = [
        ["Station", "Zone", "Lat", "Lon"],
        ["Bollaram Industrial Area", "Industrial", "17.540", "78.350"],
        ["Central University", "Residential", "17.460", "78.330"],
        ["ICRISAT Patancheru", "Rural", "17.510", "78.270"],
        ["IDA Pashamylaram", "Industrial", "17.530", "78.210"],
        ["Sanathnagar", "Commercial", "17.456", "78.443"],
        ["Zoo Park", "Residential", "17.350", "78.451"],
        ["Nacharam", "Industrial", "17.428", "78.554"],
    ]
    add_table(slide, Inches(6.8), Inches(1.6), 8, 4,
              [Inches(2.4), Inches(1.2), Inches(0.9), Inches(0.9)],
              stations, header_color=ACCENT_ORANGE)

    # Bottom: Dataset summary
    add_card(slide, Inches(0.5), Inches(4.5), Inches(12.1), Inches(2.4),
             "DATASET SUMMARY", [
                 "122,808 total hourly records (7 stations × 17,544 hours each)",
                 "17 raw columns: 6 pollutants + 7 meteorological + AQI + AQI_Category + station + zone",
                 "0 missing values across all stations — CAMS reanalysis provides complete global coverage",
                 "AQI distribution: Good (22,698) | Satisfactory (60,171) | Moderate (36,670) | Poor (3,267) | Very Poor (2)",
             ], accent=ACCENT_TEAL, body_size=12)

    # ══════════════════════════════════════════
    # SLIDE 6: SYSTEM ARCHITECTURE
    # ══════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_section_header(slide, "System Architecture & Pipeline",
                       "End-to-End Modular Design")

    # Pipeline flow cards
    pipeline_steps = [
        ("1. DATA\nCOLLECTION", "Open-Meteo AQ\n+ Weather APIs\n7 stations", ACCENT_BLUE),
        ("2. PREPROCESSING", "Imputation\nScaling\nTemporal features", ACCENT_TEAL),
        ("3. FEATURE\nENGINEERING", "36 features\nLags + Rolling\nCyclical encoding", ACCENT_ORANGE),
        ("4. MODEL\nTRAINING", "9 architectures\nOptuna HPO\nEarly stopping", ACCENT_PURPLE),
        ("5. EVALUATION\n& BENCHMARK", "RMSE, MAE, R²\nF1-score\nAQI categories", GREEN),
    ]

    for i, (title, desc, color) in enumerate(pipeline_steps):
        x = Inches(0.4 + i * 2.55)
        y = Inches(1.6)

        # Card
        card = add_shape_rect(slide, x, y, Inches(2.3), Inches(2.2), CARD_BG, color)

        # Title
        add_textbox(slide, x + Inches(0.1), y + Inches(0.15),
                    Inches(2.1), Inches(0.8),
                    title, font_size=13, color=color, bold=True,
                    alignment=PP_ALIGN.CENTER)

        # Description
        add_textbox(slide, x + Inches(0.1), y + Inches(1.0),
                    Inches(2.1), Inches(1.1),
                    desc, font_size=11, color=LIGHT_GRAY,
                    alignment=PP_ALIGN.CENTER)

        # Arrow between steps
        if i < len(pipeline_steps) - 1:
            add_textbox(slide, x + Inches(2.3), y + Inches(0.8),
                        Inches(0.3), Inches(0.4),
                        "→", font_size=22, color=MEDIUM_GRAY,
                        alignment=PP_ALIGN.CENTER)

    # Project structure
    add_card(slide, Inches(0.5), Inches(4.2), Inches(12.1), Inches(2.8),
             "PROJECT STRUCTURE — 19 PYTHON MODULES", [
                 "config.py — Central configuration (paths, features, dataclasses: DataConfig, ModelConfig, TrainConfig, OptunaConfig)",
                 "data/ — download.py (API fetching, AQI computation) | preprocessing.py (impute, scale, split) | dataset.py (PyTorch DataLoader)",
                 "models/ — model_factory.py + 9 model files: rnn.py, lstm.py, bilstm.py, lstm_attention.py, bilstm_attention.py, tcn.py, transformer.py, xgboost_model.py, lightgbm_model.py",
                 "training/ — trainer.py (PyTorch loop, AMP, gradient clipping, checkpointing) | optuna_tuner.py (unified DL+GB HPO)",
                 "evaluation/ — metrics.py (RMSE, MAE, R², MAPE, F1) | visualizer.py (loss curves, scatter plots, heatmaps, residuals)",
             ], accent=ACCENT_BLUE, body_size=11)

    # ══════════════════════════════════════════
    # SLIDE 7: FEATURE ENGINEERING
    # ══════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_section_header(slide, "Feature Engineering",
                       "36 Features across 4 Categories")

    # Feature category cards
    categories = [
        ("POLLUTANT FEATURES (6)", [
            "PM2.5 — Fine particulate matter (µg/m³)",
            "PM10 — Coarse particulate matter (µg/m³)",
            "NO₂ — Nitrogen dioxide (µg/m³)",
            "SO₂ — Sulphur dioxide (µg/m³)",
            "CO — Carbon monoxide (mg/m³)",
            "O₃ — Ground-level ozone (µg/m³)",
        ], ACCENT_BLUE),
        ("METEOROLOGICAL FEATURES (7)", [
            "AT — Ambient Temperature (°C)",
            "RH — Relative Humidity (%)",
            "WS — Wind Speed (m/s)",
            "WD — Wind Direction (degrees)",
            "RF — Rainfall (mm)",
            "SR — Solar Radiation (W/m²)",
            "BP — Barometric Pressure (hPa)",
        ], ACCENT_TEAL),
        ("TEMPORAL FEATURES (11)", [
            "hour, day_of_week, day_of_month, month",
            "Cyclical: hour_sin/cos, day_of_week_sin/cos",
            "month_sin/cos — Seasonal patterns",
            "is_weekend — Binary weekend flag",
            "Captures diurnal & seasonal AQI cycles",
            "",
        ], ACCENT_ORANGE),
        ("LAG & ROLLING FEATURES (12)", [
            "AQI_lag_1, _3, _6, _12, _24, _48",
            "AQI_rolling_mean_6, _12, _24",
            "AQI_rolling_std_6, _12, _24",
            "Autoregressive signal for temporal patterns",
            "Rolling stats capture local volatility",
            "",
        ], ACCENT_PURPLE),
    ]

    for i, (title, items, color) in enumerate(categories):
        col = i % 2
        row = i // 2
        x = Inches(0.5 + col * 6.2)
        y = Inches(1.6 + row * 2.75)
        add_card(slide, x, y, Inches(5.9), Inches(2.5),
                 title, items, accent=color, body_size=12)

    # ══════════════════════════════════════════
    # SLIDE 8: DEEP LEARNING MODELS — OVERVIEW
    # ══════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_section_header(slide, "Deep Learning Models",
                       "7 Neural Architectures for Time-Series Forecasting")

    models_data = [
        ["Model", "Architecture", "Key Mechanism", "Parameters"],
        ["RNN", "Elman RNN → FC Head", "Tanh activation, sequential processing", "~50K"],
        ["LSTM", "Stacked LSTM + LayerNorm", "Gated memory cells (forget/input/output gates)", "~200K"],
        ["BiLSTM", "Bidirectional LSTM + LN", "Forward + backward hidden state concatenation", "~350K"],
        ["LSTM+Attention", "LSTM + Bahdanau Attention", "Additive attention: V·tanh(W₁h + W₂s)", "~250K"],
        ["BiLSTM+Attention", "BiLSTM + Bahdanau Attn", "Attention over 2×hidden bidirectional states", "~400K"],
        ["TCN", "Causal Dilated Conv + Residual", "Exponential dilation (2⁰,2¹,...), weight norm", "~150K"],
        ["Transformer", "Encoder-only + Pos. Encoding", "Multi-head self-attention + causal mask + GELU", "~300K"],
    ]

    add_table(slide, Inches(0.5), Inches(1.6), 8, 4,
              [Inches(2.0), Inches(3.0), Inches(4.5), Inches(1.5)],
              models_data, header_color=ACCENT_BLUE)

    # Shared config
    add_card(slide, Inches(0.5), Inches(5.3), Inches(12.1), Inches(1.6),
             "SHARED TRAINING CONFIGURATION", [
                 "Loss: HuberLoss (δ=1.0) — robust to AQI outliers | Optimizer: AdamW (weight_decay=1e-5) | Scheduler: CosineAnnealing",
                 "Input: (batch, 72, 36) — 72-hour lookback × 36 features | Output: (batch, 24) — 24-hour AQI forecast",
                 "AMP (mixed precision) on GPU | Gradient clipping (max_norm=1.0) | Early stopping (patience=15, min_delta=1e-4)",
             ], accent=ACCENT_TEAL, body_size=12)

    # ══════════════════════════════════════════
    # SLIDE 9: RECURRENT MODELS DEEP DIVE
    # ══════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_section_header(slide, "Recurrent Architectures — Deep Dive",
                       "RNN → LSTM → BiLSTM — Progressive Complexity")

    # RNN
    add_card(slide, Inches(0.5), Inches(1.6), Inches(3.8), Inches(5.0),
             "VANILLA RNN", [
                 "Architecture:",
                 "  Input → RNN(hidden=128, layers=2)",
                 "  → Dropout(0.2)",
                 "  → FC(128→64→ReLU→24)",
                 "",
                 "Key Properties:",
                 "  • Tanh nonlinearity",
                 "  • Suffers vanishing gradients",
                 "  • Fastest to train",
                 "  • Baseline architecture",
                 "  • ~50K parameters",
             ], accent=ACCENT_BLUE, body_size=11)

    # LSTM
    add_card(slide, Inches(4.6), Inches(1.6), Inches(3.8), Inches(5.0),
             "LSTM", [
                 "Architecture:",
                 "  Input → LSTM(hidden=128, layers=2)",
                 "  → LayerNorm(128)",
                 "  → Dropout(0.2)",
                 "  → FC(128→64→ReLU→24)",
                 "",
                 "Key Properties:",
                 "  • 3 gates: forget, input, output",
                 "  • Cell state carries long-range info",
                 "  • LayerNorm stabilizes training",
                 "  • ~200K parameters",
             ], accent=ACCENT_TEAL, body_size=11)

    # BiLSTM
    add_card(slide, Inches(8.7), Inches(1.6), Inches(3.8), Inches(5.0),
             "BiLSTM", [
                 "Architecture:",
                 "  Input → BiLSTM(hidden=128, layers=2)",
                 "  → LayerNorm(256)",
                 "  → Dropout(0.2)",
                 "  → FC(256→128→ReLU→24)",
                 "",
                 "Key Properties:",
                 "  • Forward + backward passes",
                 "  • Output dim = 2 × hidden_dim",
                 "  • Captures future context",
                 "  • ~350K parameters",
             ], accent=ACCENT_ORANGE, body_size=11)

    # ══════════════════════════════════════════
    # SLIDE 10: ATTENTION MODELS DEEP DIVE
    # ══════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_section_header(slide, "Attention Mechanisms — Deep Dive",
                       "Bahdanau (Additive) Attention over Temporal States")

    # Attention mechanism explanation
    add_card(slide, Inches(0.5), Inches(1.6), Inches(6.0), Inches(3.0),
             "BAHDANAU ATTENTION MECHANISM", [
                 "Score: eᵢ = Vᵀ · tanh(W₁ · hᵢ + W₂ · s)",
                 "  hᵢ = encoder hidden states (all timesteps)",
                 "  s = decoder state (last hidden state)",
                 "",
                 "Weights: αᵢ = softmax(eᵢ)  → attention distribution",
                 "Context: c = Σ αᵢ · hᵢ  → weighted sum",
                 "",
                 "Learns WHICH of the 72 past hours matter most",
                 "Interpretable via attention weight visualization",
             ], accent=ACCENT_PURPLE, body_size=12)

    # LSTM+Attention
    add_card(slide, Inches(0.5), Inches(4.9), Inches(5.8), Inches(2.0),
             "LSTM + ATTENTION", [
                 "LSTM(128, 2 layers) → Attention(hidden=128, attn_dim=64) → LayerNorm → FC(128→64→24)",
                 "Decoder query: h_n[-1] (last layer hidden state)",
                 "Attention over 72 LSTM output states → context vector → prediction",
             ], accent=ACCENT_BLUE, body_size=12)

    # BiLSTM+Attention
    add_card(slide, Inches(6.8), Inches(1.6), Inches(5.8), Inches(3.0),
             "BiLSTM + ATTENTION", [
                 "BiLSTM(128, 2 layers) → outputs (batch, 72, 256)",
                 "Decoder query: cat(h_n[-2], h_n[-1]) = (batch, 256)",
                 "Attention over 256-dim bidirectional states",
                 "→ LayerNorm(256) → FC(256→128→ReLU→24)",
                 "",
                 "Combines BiLSTM's bidirectional context with",
                 "attention's dynamic temporal weighting",
             ], accent=ACCENT_TEAL, body_size=12)

    # get_attention_weights
    add_card(slide, Inches(6.8), Inches(4.9), Inches(5.8), Inches(2.0),
             "INTERPRETABILITY", [
                 "Both models expose get_attention_weights(x) method",
                 "Returns (batch, 72) attention distribution — visualizable as heatmap",
                 "Reveals which past hours drive the 24h forecast (e.g., rush hour spikes)",
             ], accent=ACCENT_ORANGE, body_size=12)

    # ══════════════════════════════════════════
    # SLIDE 11: TCN & TRANSFORMER
    # ══════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_section_header(slide, "TCN & Transformer — Modern Architectures",
                       "Convolutional & Self-Attention Approaches")

    # TCN
    add_card(slide, Inches(0.5), Inches(1.6), Inches(5.8), Inches(5.2),
             "TEMPORAL CONVOLUTIONAL NETWORK (TCN)", [
                 "Architecture: 5 stacked TemporalBlocks",
                 "  CausalConv1d → ReLU → Dropout → CausalConv1d → + Residual",
                 "",
                 "Key Design:",
                 "  • Causal padding: output at t depends only on t and earlier",
                 "  • Dilations: [1, 2, 4, 8, 16] — exponential growth",
                 "  • Receptive Field: 1 + 2×(3-1)×(1+2+4+8+16) = 125 > 72 ✓",
                 "  • Weight normalization on all conv layers",
                 "  • 1×1 conv for residual channel matching",
                 "",
                 "Advantages over RNNs:",
                 "  • Fully parallelizable (no sequential dependency)",
                 "  • Stable gradients (no vanishing/exploding)",
                 "  • O(log n) depth for O(n) receptive field",
                 "",
                 "Global avg pooling → FC(64→32→ReLU→24)",
             ], accent=ACCENT_ORANGE, body_size=11)

    # Transformer
    add_card(slide, Inches(6.8), Inches(1.6), Inches(5.8), Inches(5.2),
             "TRANSFORMER ENCODER", [
                 "Architecture: Encoder-only (Vaswani et al., 2017)",
                 "  Input → Linear(36→128) × √d_model",
                 "  → Sinusoidal Positional Encoding",
                 "  → 3× TransformerEncoderLayer",
                 "    (d=128, nhead=8, ff=256, GELU, Pre-LN)",
                 "",
                 "Key Design:",
                 "  • Multi-head self-attention: 8 heads × 16 dim each",
                 "  • Pre-LayerNorm (norm_first=True) — stable training",
                 "  • Causal mask: upper-triangular -inf prevents",
                 "    attending to future timesteps",
                 "  • GELU activation in feedforward sublayers",
                 "",
                 "d_model % nhead == 0 assertion enforced",
                 "Global avg pooling → FC(128→64→GELU→24)",
             ], accent=ACCENT_PURPLE, body_size=11)

    # ══════════════════════════════════════════
    # SLIDE 12: GRADIENT BOOSTING MODELS
    # ══════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_section_header(slide, "Gradient Boosting Models",
                       "XGBoost & LightGBM — Tree-Based Approaches")

    # XGBoost
    add_card(slide, Inches(0.5), Inches(1.6), Inches(5.8), Inches(3.5),
             "XGBOOST", [
                 "Histogram-based gradient boosting (Chen & Guestrin, 2016)",
                 "1000 estimators, max_depth=8, lr=0.05",
                 "Regularization: α=0.1, λ=1.0, subsample=0.8",
                 "GPU-accelerated via tree_method='hist', device='cuda'",
                 "Early stopping: 50 rounds on validation RMSE",
                 "",
                 "Operates on TABULAR features (no sliding windows)",
                 "Uses lag + rolling features as autoregressive inputs",
                 "Predicts single-step AQI (unlike DL 24-step horizon)",
             ], accent=ACCENT_BLUE, body_size=12)

    # LightGBM
    add_card(slide, Inches(6.8), Inches(1.6), Inches(5.8), Inches(3.5),
             "LIGHTGBM", [
                 "Leaf-wise tree growth (Ke et al., 2017)",
                 "1000 estimators, num_leaves=31, max_depth=8",
                 "GBDT boosting, RMSE objective, lr=0.05",
                 "GPU-accelerated when available",
                 "Callbacks API for early stopping + log evaluation",
                 "",
                 "Advantages over XGBoost:",
                 "  • Faster training with histogram binning",
                 "  • Better handling of high-cardinality features",
                 "  • Leaf-wise growth → deeper, more accurate trees",
             ], accent=ACCENT_TEAL, body_size=12)

    # Comparison
    add_card(slide, Inches(0.5), Inches(5.4), Inches(12.1), Inches(1.5),
             "DL vs GB — KEY DIFFERENCES", [
                 "DL models: Input = (72, 36) sliding window → Output = 24-step horizon | Learn temporal patterns from raw sequences",
                 "GB models: Input = 36 tabular features (with lags) → Output = 1-step prediction | Feature engineering drives performance",
                 "Both families auto-tuned by Optuna HPO with model-specific search spaces",
             ], accent=ACCENT_ORANGE, body_size=12)

    # ══════════════════════════════════════════
    # SLIDE 13: OPTUNA HPO
    # ══════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_section_header(slide, "Hyperparameter Optimization",
                       "Optuna — Automated, Efficient Search")

    # Optuna config
    add_card(slide, Inches(0.5), Inches(1.6), Inches(5.8), Inches(2.5),
             "OPTUNA CONFIGURATION", [
                 "n_trials = 50 per model (450 total trials)",
                 "timeout = 3600 seconds (1 hour/model max)",
                 "Sampler: TPE (Tree-structured Parzen Estimator)",
                 "Pruner: MedianPruner (n_startup=5, n_warmup=10)",
                 "Direction: minimize validation RMSE",
             ], accent=ACCENT_PURPLE)

    # DL search space
    add_card(slide, Inches(6.8), Inches(1.6), Inches(5.8), Inches(2.5),
             "DL MODEL SEARCH SPACE", [
                 "hidden_dim ∈ {32, 64, 128, 256}",
                 "num_layers ∈ [1, 4]",
                 "dropout ∈ [0.1, 0.5]",
                 "learning_rate ∈ [1e-5, 1e-2] (log scale)",
                 "scheduler ∈ {cosine, plateau}",
             ], accent=ACCENT_BLUE)

    # GB search space
    add_card(slide, Inches(0.5), Inches(4.4), Inches(5.8), Inches(2.5),
             "GB MODEL SEARCH SPACE", [
                 "n_estimators ∈ [100, 2000]",
                 "max_depth ∈ [3, 12]",
                 "learning_rate ∈ [0.005, 0.3] (log)",
                 "subsample ∈ [0.5, 1.0]",
                 "reg_alpha, reg_lambda ∈ [1e-3, 10.0] (log)",
             ], accent=ACCENT_TEAL)

    # Model-specific
    add_card(slide, Inches(6.8), Inches(4.4), Inches(5.8), Inches(2.5),
             "MODEL-SPECIFIC SEARCH SPACES", [
                 "TCN: channels ∈ {32,64,128}, blocks ∈ [3,6], kernel ∈ {3,5,7}",
                 "Transformer: d_model ∈ {64,128,256}, nhead ∈ {4,8}, layers ∈ [2,6]",
                 "LightGBM: num_leaves ∈ [15,127], boosting ∈ {gbdt, dart}",
                 "XGBoost: min_child_weight ∈ [1,20], gamma ∈ [0,5]",
                 "Constraint: LightGBM num_leaves < 2^max_depth",
             ], accent=ACCENT_ORANGE)

    # ══════════════════════════════════════════
    # SLIDE 14: TRAINING STRATEGY
    # ══════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_section_header(slide, "Training Strategy & Implementation",
                       "Production-Grade Training Loop")

    # Left column
    add_card(slide, Inches(0.5), Inches(1.6), Inches(3.8), Inches(3.0),
             "TRAINING LOOP", [
                 "100 epochs (max), early stopping at 15",
                 "AdamW optimizer (lr=1e-3, wd=1e-5)",
                 "CosineAnnealingLR (eta_min=1e-7)",
                 "HuberLoss (δ=1.0): robust to AQI outliers",
                 "Gradient clipping: max_norm=1.0",
                 "Best model checkpoint saving",
             ], accent=ACCENT_BLUE, body_size=12)

    # Middle column
    add_card(slide, Inches(4.6), Inches(1.6), Inches(3.8), Inches(3.0),
             "GPU ACCELERATION", [
                 "CUDA 12.1 with PyTorch 2.2.2",
                 "Automatic Mixed Precision (AMP)",
                 "GradScaler for FP16 training",
                 "XGBoost: tree_method='hist'+GPU",
                 "LightGBM: device='gpu'",
                 "config.DEVICE handles fallback",
             ], accent=ACCENT_TEAL, body_size=12)

    # Right column
    add_card(slide, Inches(8.7), Inches(1.6), Inches(3.8), Inches(3.0),
             "DATA INTEGRITY", [
                 "Chronological split: 70/15/15",
                 "NO shuffling (time-series constraint)",
                 "Scalers fit on train set ONLY",
                 "Val/test use transform-only mode",
                 "Negative value clipping post-imputation",
                 "All-NaN column guard in preprocessing",
             ], accent=ACCENT_ORANGE, body_size=12)

    # Reproducibility
    add_card(slide, Inches(0.5), Inches(5.0), Inches(12.1), Inches(1.8),
             "REPRODUCIBILITY & SAFEGUARDS", [
                 "Seed control: np.random.seed(42), torch.manual_seed(42), torch.cuda.manual_seed_all(42), cudnn.deterministic=True",
                 "Checkpointing: Best model saved to checkpoints/{model_name}_best.pt with optimizer state + scheduler state + epoch",
                 "Logging: dual output (stdout + results/pipeline.log) with structured timestamps for debugging",
                 "Inverse scaling: predictions + targets inverse-transformed BEFORE metric computation (prevents z-score RMSE/MAE)",
             ], accent=ACCENT_PURPLE, body_size=12)

    # ══════════════════════════════════════════
    # SLIDE 15: EVALUATION METRICS
    # ══════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_section_header(slide, "Evaluation Metrics & Methodology",
                       "Comprehensive Regression + Classification Assessment")

    # Regression metrics
    add_card(slide, Inches(0.5), Inches(1.6), Inches(5.8), Inches(3.0),
             "REGRESSION METRICS", [
                 "RMSE — Root Mean Squared Error: √(Σ(yᵢ-ŷᵢ)²/n)",
                 "    Penalizes large errors; primary ranking metric",
                 "MAE — Mean Absolute Error: Σ|yᵢ-ŷᵢ|/n",
                 "    Robust to outliers; interpretable in AQI units",
                 "R² — Coefficient of Determination: 1 - SS_res/SS_tot",
                 "    Proportion of variance explained (target > 0.85)",
                 "MAPE — Mean Absolute Percentage Error",
                 "    Scale-independent accuracy measure (%)",
             ], accent=ACCENT_BLUE, body_size=12)

    # Classification metrics
    add_card(slide, Inches(6.8), Inches(1.6), Inches(5.8), Inches(3.0),
             "CLASSIFICATION METRICS", [
                 "AQI values mapped to 6 NAQI categories:",
                 "  Good (0-50) | Satisfactory (51-100) | Moderate (101-200)",
                 "  Poor (201-300) | Very Poor (301-400) | Severe (401-500)",
                 "",
                 "F1_weighted — Weighted by class frequency",
                 "F1_macro — Unweighted average across categories",
                 "Per-class F1 — Individual category accuracy",
                 "Confusion matrix for detailed error analysis",
             ], accent=ACCENT_TEAL, body_size=12)

    # Methodology
    add_card(slide, Inches(0.5), Inches(4.9), Inches(12.1), Inches(2.0),
             "EVALUATION METHODOLOGY", [
                 "All metrics computed on held-out TEST SET only (final 15% = 2,625 hours ≈ 109 days, Sept–Dec 2024)",
                 "Predictions inverse-scaled from z-scores back to original AQI scale before metric computation",
                 "Complete metric suite: {RMSE, MAE, R², MAPE, F1_weighted, F1_macro} enables multi-dimensional model comparison",
                 "Visualization: loss curves, actual-vs-predicted line/scatter, metric heatmap, residual distributions, feature importance",
             ], accent=ACCENT_PURPLE, body_size=12)

    # ══════════════════════════════════════════
    # SLIDE 16: DATA ANALYSIS
    # ══════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_section_header(slide, "Data Analysis & Statistics",
                       "122,808 Hourly Records — 2 Years of Hyderabad Air Quality")

    # Dataset overview
    stats_data = [
        ["Metric", "Value"],
        ["Total Records", "122,808 hourly observations"],
        ["Stations", "7 CPCB CAAQMS sites across Hyderabad"],
        ["Date Range", "Jan 1, 2023 — Dec 31, 2024 (731 days)"],
        ["Missing Values", "0 (CAMS reanalysis provides complete coverage)"],
        ["Features (raw)", "17 columns (6 pollutant + 7 meteo + 4 metadata)"],
        ["Features (engineered)", "36 (after temporal encoding + lags + rolling)"],
        ["Train Split", "12,247 hours (Jan 2023 — May 2024)"],
        ["Validation Split", "2,624 hours (May — Sep 2024)"],
        ["Test Split", "2,625 hours (Sep — Dec 2024)"],
    ]

    add_table(slide, Inches(0.5), Inches(1.6), len(stats_data), 2,
              [Inches(2.5), Inches(5.5)], stats_data, header_color=ACCENT_BLUE)

    # AQI distribution
    aqi_data = [
        ["AQI Category", "Range", "Count", "Percentage", "Health Impact"],
        ["Good", "0 – 50", "22,698", "18.5%", "Minimal impact"],
        ["Satisfactory", "51 – 100", "60,171", "49.0%", "Minor breathing discomfort for sensitive people"],
        ["Moderate", "101 – 200", "36,670", "29.9%", "Breathing discomfort for asthmatics"],
        ["Poor", "201 – 300", "3,267", "2.7%", "Breathing discomfort for most people"],
        ["Very Poor", "301 – 400", "2", "<0.01%", "Respiratory illness on prolonged exposure"],
        ["Severe", "401 – 500", "0", "0%", "Health emergency; affects healthy people"],
    ]

    add_table(slide, Inches(0.5), Inches(4.6), len(aqi_data), 5,
              [Inches(1.5), Inches(1.0), Inches(1.2), Inches(1.2), Inches(4.0)],
              aqi_data, header_color=ACCENT_TEAL)

    # ══════════════════════════════════════════
    # SLIDE 17: TECH STACK
    # ══════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_section_header(slide, "Technology Stack",
                       "Frameworks, Libraries & Infrastructure")

    tech_categories = [
        ("CORE ML/DL", [
            "Python 3.10.13",
            "PyTorch 2.2.2 + CUDA 12.1",
            "scikit-learn 1.5.2",
            "NumPy 1.26.4",
            "Pandas 2.3.3",
        ], ACCENT_BLUE),
        ("GRADIENT BOOSTING", [
            "XGBoost 2.0.3 (GPU)",
            "LightGBM 4.6.0 (GPU)",
            "GPU tree_method='hist'",
            "Leaf-wise growth strategy",
            "Early stopping callbacks",
        ], ACCENT_TEAL),
        ("OPTIMIZATION", [
            "Optuna 4.7.0",
            "TPE Sampler",
            "MedianPruner",
            "50 trials × 9 models",
            "Parallel trial evaluation",
        ], ACCENT_ORANGE),
        ("VISUALIZATION & INFRA", [
            "Matplotlib 3.9.1",
            "Seaborn 0.13.2",
            "Conda (dl-env)",
            "VS Code + GitHub Copilot",
            "Multi-agent review system",
        ], ACCENT_PURPLE),
    ]

    for i, (title, items, color) in enumerate(tech_categories):
        x = Inches(0.4 + i * 3.15)
        add_card(slide, x, Inches(1.6), Inches(2.9), Inches(3.0),
                 title, items, accent=color, body_size=12)

    # Hardware
    add_card(slide, Inches(0.5), Inches(5.0), Inches(12.1), Inches(1.8),
             "HARDWARE & ENVIRONMENT", [
                 "GPU: NVIDIA CUDA 12.1 compatible | CPU fallback via config.DEVICE (torch.device auto-detection)",
                 "Environment: Conda virtual environment 'dl-env' | Python 3.10.13 | Windows",
                 "Storage: ~200MB raw data (7 station CSVs) | ~50MB processed pipeline output (pickle) | ~500MB model checkpoints",
             ], accent=GREEN, body_size=12)

    # ══════════════════════════════════════════
    # SLIDE 18: MULTI-AGENT SYSTEM
    # ══════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_section_header(slide, "Multi-Agent Development System",
                       "11 Specialized AI Agents for Code Quality")

    # Agent architecture
    agents_data = [
        ["Agent", "Scope", "Responsibility"],
        ["rnn-specialist", "models/rnn.py", "RNN architecture validation"],
        ["lstm-specialist", "models/lstm.py", "LSTM gate mechanics, LayerNorm"],
        ["bilstm-specialist", "models/bilstm.py", "BiLSTM dimension handling"],
        ["lstm-attention-specialist", "models/lstm_attention.py", "Attention weight computation"],
        ["bilstm-attention-specialist", "models/bilstm_attention.py", "BiLSTM + Attention integration"],
        ["tcn-specialist", "models/tcn.py", "Causal convolution, receptive field"],
        ["transformer-specialist", "models/transformer.py", "Multi-head attention, causal mask"],
        ["xgboost-specialist", "models/xgboost_model.py", "GPU config, feature importance"],
        ["lightgbm-specialist", "models/lightgbm_model.py", "Callbacks API, leaf-wise growth"],
        ["orchestrator", "main.py + pipeline", "Coordinate training, benchmarking"],
        ["quality-monitor", "All modules", "Data leakage, convergence validation"],
    ]

    add_table(slide, Inches(0.5), Inches(1.6), len(agents_data), 3,
              [Inches(2.8), Inches(3.0), Inches(4.0)],
              agents_data, header_color=ACCENT_PURPLE)

    # Bugs found
    add_card(slide, Inches(0.5), Inches(5.6), Inches(12.1), Inches(1.4),
             "CRITICAL BUGS FOUND & FIXED BY AGENTS", [
                 "KeyError 'feature_cols' → 'feature_names' | y_pred/y_true swap in metrics | shuffle=True on train DataLoader (time-series violation)",
                 "No inverse scaling before metrics (z-score RMSE) | All-NaN column propagation | Negative values from interpolation | Metric key case mismatch",
             ], accent=RED, body_size=12)

    # ══════════════════════════════════════════
    # SLIDE 19: RESULTS & FUTURE WORK
    # ══════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_section_header(slide, "Results & Future Work",
                       "Current Status and Next Steps")

    # Current status
    add_card(slide, Inches(0.5), Inches(1.6), Inches(5.8), Inches(3.0),
             "CURRENT STATUS ✓", [
                 "✓ 122,808 hourly records collected from 7 stations",
                 "✓ Full preprocessing pipeline validated (36 features)",
                 "✓ 19 Python modules implemented and reviewed",
                 "✓ 11 AI agents validated all model architectures",
                 "✓ 8 critical bugs found and fixed pre-training",
                 "✓ Pipeline ready for training execution",
             ], accent=GREEN, body_size=13)

    # Future work
    add_card(slide, Inches(6.8), Inches(1.6), Inches(5.8), Inches(3.0),
             "FUTURE WORK & NEXT STEPS", [
                 "→ Execute full model training across all 9 architectures",
                 "→ Run Optuna HPO (50 trials/model, ~450 total trials)",
                 "→ Generate comprehensive benchmark comparison",
                 "→ Ensemble methods: stacking top-3 models",
                 "→ Real-time prediction API deployment (FastAPI)",
                 "→ Mobile app with push notifications for AQI alerts",
             ], accent=ACCENT_BLUE, body_size=13)

    # Expected outcomes
    add_card(slide, Inches(0.5), Inches(4.9), Inches(12.1), Inches(2.0),
             "EXPECTED OUTCOMES", [
                 "Identify the best-performing architecture for Hyderabad AQI prediction — hypothesis: BiLSTM+Attention or Transformer will lead DL; XGBoost competitive for tabular features",
                 "Achieve R² > 0.85, RMSE < 15 AQI points, F1_weighted > 0.80 for 24-hour forecasts across all seasons",
                 "Publishable benchmark: first 9-model comparison on Hyderabad hourly AQ data with Optuna-tuned hyperparameters",
                 "Open-source, reproducible pipeline for other Indian cities (Delhi, Mumbai, Chennai) — only requires changing station coordinates",
             ], accent=ACCENT_TEAL, body_size=12)

    # ══════════════════════════════════════════
    # SLIDE 20: THANK YOU
    # ══════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    # Decorative bars
    for i, c in enumerate([ACCENT_BLUE, ACCENT_TEAL, ACCENT_PURPLE, ACCENT_ORANGE]):
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(i * 3.33), Inches(0), Inches(3.34), Pt(5)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = c
        bar.line.fill.background()

    add_textbox(slide, Inches(0.8), Inches(2.0), Inches(11.5), Inches(1.0),
                "Thank You",
                font_size=48, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(0.8), Inches(3.2), Inches(11.5), Inches(0.6),
                "Air Quality Index Prediction System — Hyderabad, Telangana",
                font_size=18, color=ACCENT_TEAL, alignment=PP_ALIGN.CENTER)

    # Divider
    div = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(5.0), Inches(4.0), Inches(3.33), Pt(2)
    )
    div.fill.solid()
    div.fill.fore_color.rgb = ACCENT_BLUE
    div.line.fill.background()

    # Team
    add_textbox(slide, Inches(0.8), Inches(4.4), Inches(11.5), Inches(0.5),
                "Sakshith Srihari  •  Kamal Bura  •  Abdul Muttalib",
                font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(0.8), Inches(4.9), Inches(11.5), Inches(0.4),
                "1602-22-748-011  •  1602-22-748-302  •  1602-22-748-046",
                font_size=13, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(0.8), Inches(5.6), Inches(11.5), Inches(0.5),
                "Questions & Discussion",
                font_size=20, color=ACCENT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)

    # Bottom footer
    for i, c in enumerate([ACCENT_BLUE, ACCENT_TEAL, ACCENT_PURPLE, ACCENT_ORANGE]):
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(i * 3.33), Inches(7.4), Inches(3.34), Pt(5)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = c
        bar.line.fill.background()

    return prs


# ═══════════════════════════════════════════════════════════════
if __name__ == "__main__":
    output_dir = os.path.dirname(os.path.abspath(__file__))
    project_root = os.path.dirname(output_dir)
    output_path = os.path.join(project_root, "AQI_Prediction_Presentation.pptx")

    print("Generating presentation...")
    prs = create_presentation()
    prs.save(output_path)
    print(f"Saved to: {output_path}")
    print(f"Slides: {len(prs.slides)}")
